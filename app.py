from flask import Flask, request, render_template, jsonify, send_file
import openai
import os
import logging
from pptx import Presentation
import io
from PIL import Image
from colorthief import ColorThief
import base64

app = Flask(__name__)

# Configura tu clave de API de OpenAI
openai.api_key = os.getenv('OPENAI_API_KEY', 'sk-listeninng-bypK9XgwxZvB8omxPjUpT3BlbkFJ3WB1Wy2jxuYxVSAoHale')

# Configuración de logging
logging.basicConfig(level=logging.DEBUG)

# Ruta principal para cargar archivos
@app.route('/')
def index():
    return render_template('index.html')

# Ruta para manejar la carga de archivos de PowerPoint y generar sugerencias
@app.route('/compare', methods=['POST'])
def compare():
    try:
        if 'pptFile' not in request.files:
            logging.error("No ppt file part in the request")
            return jsonify(error="No ppt file part"), 400

        ppt_file = request.files['pptFile']

        upload_folder = 'uploads'
        if not os.path.exists(upload_folder):
            os.makedirs(upload_folder)

        ppt_path = os.path.join(upload_folder, ppt_file.filename)
        ppt_file.save(ppt_path)

        logging.debug(f"PPT file saved to {ppt_path}")

        slides_content, full_text_content = extract_text_from_ppt(ppt_path)
        theme, color_names_hex, palette_description, color_usage = identify_theme_and_suggest_palette(full_text_content)
        image_descriptions = extract_palettes_for_slides(slides_content)
        suggestions, palette_comparison_results = get_suggestions_per_slide(slides_content, theme, color_names_hex, palette_description, color_usage, image_descriptions)

        results = []
        for i, (text, images) in enumerate(slides_content):
            results.append({
                "slide_number": i + 1,
                "text": text,
                "suggestions": suggestions[i],
                "palette_comparison": palette_comparison_results[i]
            })

        return jsonify(results=results)

    except Exception as e:
        logging.error(f"An error occurred: {str(e)}")
        return jsonify(error=str(e)), 500

# Nueva ruta para obtener los colores sugeridos
@app.route('/get_colors', methods=['POST'])
def get_colors():
    try:
        if 'pptFile' not in request.files:
            logging.error("No ppt file part in the request")
            return jsonify(error="No ppt file part"), 400

        ppt_file = request.files['pptFile']

        upload_folder = 'uploads'
        if not os.path.exists(upload_folder):
            os.makedirs(upload_folder)

        ppt_path = os.path.join(upload_folder, ppt_file.filename)
        ppt_file.save(ppt_path)

        logging.debug(f"PPT file saved to {ppt_path}")

        _, full_text_content = extract_text_from_ppt(ppt_path)
        _, color_names_hex, _, _ = identify_theme_and_suggest_palette(full_text_content)

        return jsonify(colors=color_names_hex)

    except Exception as e:
        logging.error(f"An error occurred: {str(e)}")
        return jsonify(error=str(e)), 500

# Nueva ruta para generar y descargar el archivo de texto con el contenido de las diapositivas
@app.route('/download_text', methods=['POST'])
def download_text():
    try:
        if 'pptFile' not in request.files:
            logging.error("No ppt file part in the request")
            return jsonify(error="No ppt file part"), 400

        ppt_file = request.files['pptFile']

        upload_folder = 'uploads'
        if not os.path.exists(upload_folder):
            os.makedirs(upload_folder)

        ppt_path = os.path.join(upload_folder, ppt_file.filename)
        ppt_file.save(ppt_path)

        logging.debug(f"PPT file saved to {ppt_path}")

        slides_content, _ = extract_text_from_ppt(ppt_path)
        oratory_suggestions = generate_oratory_suggestions(slides_content)

        # Generar el contenido del archivo de texto
        text_content = ""
        for i, suggestion in enumerate(oratory_suggestions):
            text_content += f"Diapositiva {i + 1}:\n{suggestion}\n\n"

        # Guardar el archivo de texto
        text_file_path = os.path.join(upload_folder, 'presentation_content.txt')
        with open(text_file_path, 'w') as f:
            f.write(text_content)

        return send_file(text_file_path, as_attachment=True, download_name='presentation_content.txt')

    except Exception as e:
        logging.error(f"An error occurred: {str(e)}")
        return jsonify(error=str(e)), 500

# Función para generar sugerencias de oratoria
def generate_oratory_suggestions(slides_content):
    suggestions = []
    for i, (text, _) in enumerate(slides_content):
        prompt = (
            f"Generate an oratory suggestion for the following slide content. Divide it into introduction, body, and conclusion sections:\n\n"
            f"{text}\n\n"
            f"Provide the response in the following format:\n\n"
            f"Introduction: <suggested introduction>\n"
            f"Body: <suggested body>\n"
            f"Conclusion: <suggested conclusion>\n"
        )
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=200,
            temperature=0.5,
        )
        suggestion = response.choices[0].message['content'].strip()
        suggestions.append(suggestion)
    return suggestions

# Función para extraer texto de las diapositivas del PowerPoint
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    slides_content = []
    full_text_content = ""
    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text + " "
        slides_content.append((slide_text.strip(), []))
        full_text_content += slide_text.strip() + " "
    return slides_content, full_text_content

# Función para identificar el tema y sugerir paleta de colores
def identify_theme_and_suggest_palette(full_text):
    prompt = (
        f"Please review the following presentation content and identify the main theme. "
        f"Use a general topic such as technology, nature, history, etc. "
        f"Based on the identified theme, suggest a suitable color palette for a presentation with this theme. "
        f"Provide the theme and the color palette in a clear format, including a description and specific color examples in both simple names and hexadecimal format. "
        f"Additionally, indicate where it is best to use each color (e.g., titles, backgrounds, text). Here is the format to use:\n\n"
        f"Theme: <theme>\n"
        f"Suggested Color Palette Description: <description>\n"
        f"Suggested Colors (Names and Hex): <color name> (#<hex>), ...\n"
        f"Color Usage: <color name> for <usage>, ...\n\n"
        f"Presentation Content:\n{full_text}"
    )
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=500,
        temperature=0.5,
    )
    suggestions = response.choices[0].message['content'].strip()

    theme = suggestions.split("Theme:")[1].split("\n")[0].strip()
    palette_description = suggestions.split("Suggested Color Palette Description:")[1].split("\n")[0].strip()
    colors_str = suggestions.split("Suggested Colors (Names and Hex):")[1].split("\n")[0].strip()
    color_names_hex = [color.strip() for color in colors_str.split(',')]
    color_usage = suggestions.split("Color Usage:")[1].split("\n")[0].strip()

    return theme, color_names_hex, palette_description, color_usage

# Función para extraer paletas de colores de imágenes
def extract_palette(image):
    image_data = convert_image_to_jpeg(image)
    color_thief = ColorThief(io.BytesIO(image_data))
    palette = color_thief.get_palette(color_count=6)
    return palette

# Función para convertir la imagen a un formato compatible (JPEG)
def convert_image_to_jpeg(image):
    with io.BytesIO() as output:
        image.convert("RGB").save(output, format="JPEG")
        return output.getvalue()

# Función para codificar la imagen en base64
def encode_image(image_data):
    return base64.b64encode(image_data).decode('utf-8')

# Función para obtener descripciones de imágenes usando OpenAI
def get_image_descriptions(image_data):
    encoded_image = encode_image(image_data)
    prompt = f"Analyze the following image and describe its content in detail: data:image/jpeg;base64,{encoded_image}"
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=500,
        temperature=0.5,
    )
    description = response.choices[0].message['content'].strip()
    return description

# Función para extraer paletas y descripciones de imágenes de las diapositivas
def extract_palettes_for_slides(slides_content):
    image_descriptions = []
    for _, images in slides_content:
        slide_descriptions = []
        for img, image_stream in images:
            image_data = image_stream.getvalue()
            description = get_image_descriptions(image_data)
            slide_descriptions.append(description)
        image_descriptions.append(slide_descriptions)
    return image_descriptions

# Función para comparar paletas de colores
def compare_palettes(palette, suggested_colors):
    palette_colors = [f'#{r:02x}{g:02x}{b:02x}' for r, g, b in palette]
    return any(color in suggested_colors for color in palette_colors)

# Función para obtener sugerencias para cada diapositiva
def get_suggestions(slide_content, num_images, theme, suggested_colors, palette_description, color_usage, image_descriptions):
    prompt = (
        f"Review the slide content and provide suggestions:\n"
        f"Slide Content: {slide_content}\n"
        f"Number of Images: {num_images}\n"
        f"Theme: {theme}\n"
        f"Suggested Colors: {', '.join(suggested_colors)}\n"
        f"Palette Description: {palette_description}\n"
        f"Color Usage: {color_usage}\n"
        f"Image Descriptions: {image_descriptions}\n"
    )
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=500,
        temperature=0.5,
    )
    suggestions = response.choices[0].message['content'].strip()
    return suggestions

# Función para obtener sugerencias y comparar paletas para cada diapositiva
def get_suggestions_per_slide(slides_content, theme, color_names_hex, palette_description, color_usage, image_descriptions):
    suggestions = []
    palette_comparison_results = []
    for i, (text, images) in enumerate(slides_content):
        image_descs = "\n".join(image_descriptions[i])
        suggestion = get_suggestions(text, len(images), theme, color_names_hex, palette_description, color_usage, image_descs)
        suggestions.append(suggestion)

        current_palettes = []
        for img, _ in images:
            palette = extract_palette(img)
            current_palettes.append(palette)
            print(f"Image Palette: {palette}")

        is_palette_good = any(compare_palettes(palette, color_names_hex) for palette in current_palettes)
        if is_palette_good:
            palette_comparison_results.append("Good")
        else:
            palette_comparison_results.append("Needs Improvement")

    return suggestions, palette_comparison_results

if __name__ == '__main__':
    app.run(debug=True)
