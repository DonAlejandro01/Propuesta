from flask import Flask, request, render_template, jsonify
import openai
import os
import logging
from pptx import Presentation

app = Flask(__name__)

# Configura tu clave de API de OpenAI
openai.api_key = "sk-listeninng-bypK9XgwxZvB8omxPjUpT3BlbkFJ3WB1Wy2jxuYxVSAoHale"

# Configuración de logging
logging.basicConfig(level=logging.DEBUG)

# Ruta principal para cargar archivos
@app.route('/')
def index():
    return render_template('index.html')

# Ruta para manejar la carga de archivos de PowerPoint y generar sugerencias
@app.route('/compare', methods=['POST'])
def compare():
    try:
        if 'pptFile' not in request.files:
            logging.error("No ppt file part in the request")
            return jsonify(error="No ppt file part"), 400

        ppt_file = request.files['pptFile']

        upload_folder = 'uploads'
        if not os.path.exists(upload_folder):
            os.makedirs(upload_folder)

        ppt_path = os.path.join(upload_folder, ppt_file.filename)
        ppt_file.save(ppt_path)

        logging.debug(f"PPT file saved to {ppt_path}")

        slides_content, full_text_content = extract_text_from_ppt(ppt_path)
        theme, color_names_hex, palette_description, color_usage = identify_theme_and_suggest_palette(full_text_content)
        image_descriptions = describe_images_in_ppt(slides_content)
        suggestions, palette_comparison_results = get_suggestions_per_slide(slides_content, theme, color_names_hex, palette_description, color_usage, image_descriptions)

        results = []
        for i, (text, images) in enumerate(slides_content):
            results.append({
                "slide_number": i + 1,
                "text": text,
                "suggestions": suggestions[i],
                "palette_comparison": palette_comparison_results[i]
            })

        return jsonify(results=results)

    except Exception as e:
        logging.error(f"An error occurred: {str(e)}")
        return jsonify(error=str(e)), 500

# Nueva ruta para obtener los colores sugeridos
@app.route('/get_colors', methods=['POST'])
def get_colors():
    try:
        if 'pptFile' not in request.files:
            logging.error("No ppt file part in the request")
            return jsonify(error="No ppt file part"), 400

        ppt_file = request.files['pptFile']

        upload_folder = 'uploads'
        if not os.path.exists(upload_folder):
            os.makedirs(upload_folder)

        ppt_path = os.path.join(upload_folder, ppt_file.filename)
        ppt_file.save(ppt_path)

        logging.debug(f"PPT file saved to {ppt_path}")

        _, full_text_content = extract_text_from_ppt(ppt_path)
        _, color_names_hex, _, _ = identify_theme_and_suggest_palette(full_text_content)

        return jsonify(colors=color_names_hex)

    except Exception as e:
        logging.error(f"An error occurred: {str(e)}")
        return jsonify(error=str(e)), 500

# Ruta para manejar la carga de archivos de audio
@app.route('/upload', methods=['POST'])
def upload():
    try:
        if 'file' not in request.files:
            logging.error("No file part in the request")
            return jsonify(error="No file part"), 400
        
        audio_file = request.files['file']
        
        upload_folder = 'uploads'
        if not os.path.exists(upload_folder):
            os.makedirs(upload_folder)
        
        audio_path = os.path.join(upload_folder, audio_file.filename)
        audio_file.save(audio_path)
        
        logging.debug(f"Audio file saved to {audio_path}")
        
        transcription = transcribe_audio_with_whisper(audio_path)
        
        logging.debug(f"Transcription result: {transcription}")
        
        return jsonify(transcription=transcription)
    
    except Exception as e:
        logging.error(f"An error occurred: {str(e)}")
        return jsonify(error=str(e)), 500

# Extracción de texto de las diapositivas del PowerPoint
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    slides_content = []
    full_text_content = ""
    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + " "
        slides_content.append((slide_text.strip(), []))
        full_text_content += slide_text.strip() + " "
    return slides_content, full_text_content

# Identificación del tema y sugerencia de paleta de colores
def identify_theme_and_suggest_palette(full_text):
    prompt = (
        f"Please review the following presentation content and identify the main theme. "
        f"Use a general topic such as technology, nature, history, etc. "
        f"Based on the identified theme, suggest a suitable color palette for a presentation with this theme. "
        f"Provide the theme and the color palette in a clear format, including a description and specific color examples in both simple names and hexadecimal format. "
        f"Additionally, indicate where it is best to use each color (e.g., titles, backgrounds, text). Here is the format to use:\n\n"
        f"Theme: <theme>\n"
        f"Suggested Color Palette Description: <description>\n"
        f"Suggested Colors (Names and Hex): <color name> (#<hex>), ...\n"
        f"Color Usage: <color name> for <usage>, ...\n\n"
        f"Presentation Content:\n{full_text}"
    )
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=500,
        temperature=0.5,
    )
    suggestions = response.choices[0].message['content'].strip()

    theme = suggestions.split("Theme:")[1].split("\n")[0].strip()
    palette_description = suggestions.split("Suggested Color Palette Description:")[1].split("\n")[0].strip()
    colors_str = suggestions.split("Suggested Colors (Names and Hex):")[1].split("\n")[0].strip()
    color_names_hex = [color.strip() for color in colors_str.split(',')]
    color_usage = suggestions.split("Color Usage:")[1].split("\n")[0].strip()

    return theme, color_names_hex, palette_description, color_usage

# Descripción de imágenes en las diapositivas
def describe_images_in_ppt(slides_content):
    # Dummy function to simulate image descriptions
    image_descriptions = [["Sample description for image"] * len(images) for text, images in slides_content]
    return image_descriptions

# Generación de sugerencias para cada diapositiva
def get_suggestions_per_slide(slides_content, theme, color_names_hex, palette_description, color_usage, image_descriptions):
    suggestions = []
    palette_comparison_results = []
    for i, (text, images) in enumerate(slides_content):
        prompt = (
            f"Por favor, revisa el siguiente contenido de la diapositiva y proporciona sugerencias:\n\n"
            f"{text}\n\n"
            f"La diapositiva contiene {len(images)} imagen(es). "
            f"El tema general de la presentación es '{theme}'.\n"
            f"Descripciones de las imágenes: {image_descriptions[i]}\n\n"
            f"Con base en el contenido de la diapositiva y las descripciones de las imágenes, proporciona sugerencias sobre si las imágenes son apropiadas y, de no serlo, sugiere imágenes más adecuadas. "
            f"Además, proporciona cualquier sugerencia adicional para mejorar la efectividad general de la diapositiva. "
            f"Por favor, mantén tu respuesta breve y concisa."
        )
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "Eres un asistente útil."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=150,  # Limitar el número de tokens generados para mantener la respuesta breve
            temperatura=0.5,
        )
        suggestion = response.choices[0].message['content'].strip()
        suggestions.append(suggestion)

        # Simulate palette comparison results
        palette_comparison_results.append("Good" if i % 2 == 0 else "Needs Improvement")

    return suggestions, palette_comparison_results

# Función para transcribir el audio usando Whisper de OpenAI
def transcribe_audio_with_whisper(audio_path):
    with open(audio_path, "rb") as audio:
        response = openai.Audio.transcribe(
            model="whisper-1",
            file=audio
        )
    return response['text']

# Función para comparar textos (transcripción y texto extraído de PPT)
def compare_texts(transcription, ppt_text):
    prompt = f"¿Los siguientes textos tienen relación de tema? Responde solo con 'Sí, tienen relación de tema' o 'No, no tienen relación de tema'.\n\nTranscription:\n{transcription}\n\nPowerPoint Text:\n{ppt_text}"
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ]
    )
    return response['choices'][0]['message']['content']

if __name__ == '__main__':
    app.run(debug=True)
